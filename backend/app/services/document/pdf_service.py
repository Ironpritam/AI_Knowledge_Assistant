from pathlib import Path

import fitz
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from unstructured.partition.auto import partition

from app.services.document.text_cleaner import TextCleaner


class PDFService:
    @staticmethod
    def _extract_text_from_unstructured(file_path: Path) -> str:
        elements = partition(filename=str(file_path))
        text_parts = []

        for element in elements:
            text = getattr(element, "text", None)
            if text:
                text_parts.append(text)

        if text_parts:
            return "\n\n".join(text_parts)

        return file_path.read_text(encoding="utf-8", errors="replace")

    @staticmethod
    def _extract_docx_text(file_path: Path) -> str:
        document = DocxDocument(str(file_path))
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        return "\n\n".join(paragraphs)

    @staticmethod
    def _extract_pptx_text(file_path: Path) -> str:
        presentation = Presentation(str(file_path))
        slide_text = []

        for slide in presentation.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text.strip())
            if texts:
                slide_text.append("\n".join(texts))

        return "\n\n".join(slide_text)

    @staticmethod
    def _extract_xlsx_text(file_path: Path) -> str:
        workbook = load_workbook(file_path, read_only=True, data_only=True)
        sheet_text = []

        for worksheet in workbook.worksheets:
            rows = []
            for row in worksheet.iter_rows(values_only=True):
                row_values = ["" if value is None else str(value) for value in row]
                if any(row_values):
                    rows.append("\t".join(row_values))
            if rows:
                sheet_text.append("\n".join(rows))

        return "\n\n".join(sheet_text)

    @staticmethod
    def extract_document(file_path: Path) -> dict:
        """
        Extract text and metadata from PDFs and supported office/text/code documents.
        """
        extension = file_path.suffix.lower()

        if extension == ".pdf":
            try:
                document = fitz.open(file_path)
            except Exception as exc:
                raise RuntimeError(f"Unable to process PDF: {exc}") from exc

            pages = []

            for page_number, page in enumerate(document, start=1):
                pages.append(
                    {
                        "page": page_number,
                        "text": TextCleaner.clean(page.get_text("text")),
                    }
                )

            result = {
                "filename": file_path.name,
                "page_count": len(document),
                "metadata": document.metadata,
                "pages": pages,
            }

            document.close()
            return result

        try:
            if extension in {".docx"}:
                text = PDFService._extract_docx_text(file_path)
            elif extension in {".pptx"}:
                text = PDFService._extract_pptx_text(file_path)
            elif extension in {".xlsx"}:
                text = PDFService._extract_xlsx_text(file_path)
            elif extension in {".doc", ".ppt", ".xls"}:
                text = PDFService._extract_text_from_unstructured(file_path)
            elif extension in {".txt", ".md", ".csv", ".log", ".json", ".py", ".js", ".ts", ".tsx",
                                ".jsx", ".java", ".html", ".htm", ".css", ".sql", ".xml", ".yaml", ".yml"}:
                text = file_path.read_text(encoding="utf-8", errors="replace")
            else:
                raise ValueError(f"Unsupported file type: {extension}")
        except Exception as exc:
            raise RuntimeError(f"Unable to process document '{file_path.name}': {exc}") from exc

        cleaned_text = TextCleaner.clean(text)
        return {
            "filename": file_path.name,
            "page_count": 1,
            "metadata": {"source_type": extension.lstrip("."), "file_type": extension},
            "pages": [{"page": 1, "text": cleaned_text}],
        }