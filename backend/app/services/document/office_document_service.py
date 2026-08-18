from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

from app.services.document.base_document_service import BaseDocumentService
from app.services.document.text_cleaner import TextCleaner


class OfficeDocumentService(BaseDocumentService):
    @staticmethod
    def _extract_docx_text(file_path: Path) -> str:
        document = DocxDocument(str(file_path))
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        return "\n\n".join(paragraphs)

    @staticmethod
    def _extract_pptx_text(file_path: Path) -> str:
        presentation = Presentation(str(file_path))
        slide_text = []

        for slide in presentation.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text.strip())
            if texts:
                slide_text.append("\n".join(texts))

        return "\n\n".join(slide_text)

    @staticmethod
    def _extract_xlsx_text(file_path: Path) -> str:
        workbook = load_workbook(file_path, read_only=True, data_only=True)
        sheet_text = []

        for worksheet in workbook.worksheets:
            rows = []
            for row in worksheet.iter_rows(values_only=True):
                row_values = ["" if value is None else str(value) for value in row]
                if any(row_values):
                    rows.append("\t".join(row_values))
            if rows:
                sheet_text.append("\n".join(rows))

        return "\n\n".join(sheet_text)

    @staticmethod
    def extract_document(file_path: Path) -> dict:
        extension = file_path.suffix.lower()

        if extension == ".docx":
            text = OfficeDocumentService._extract_docx_text(file_path)
        elif extension == ".pptx":
            text = OfficeDocumentService._extract_pptx_text(file_path)
        elif extension == ".xlsx":
            text = OfficeDocumentService._extract_xlsx_text(file_path)
        else:
            raise ValueError(f"Unsupported office format: {extension}")

        cleaned = TextCleaner.clean(text)
        return {
            "filename": file_path.name,
            "page_count": 1,
            "metadata": {"source_type": extension.lstrip("."), "file_type": extension},
            "pages": [{"page": 1, "text": cleaned}],
        }
